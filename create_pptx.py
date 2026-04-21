from pptx import Presentation
import os

# Initialize presentation
prs = Presentation()

# Helper function to create bullet point slides
def add_bullet_slide(prs, title_text, bullets):
    # layout 1 is typically title and content (bullet points)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = title_text
    
    tf = slide.placeholders[1].text_frame
    for i, bullet in enumerate(bullets):
        if i == 0:
            # Overwrite the first existing paragraph
            p = tf.paragraphs[0]
        else:
            # Add new paragraphs for remaining bullets
            p = tf.add_paragraph()
        p.text = bullet
        p.level = 0

# Slide 1: Title Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
title_1 = slide_1.shapes.title
subtitle_1 = slide_1.placeholders[1]
title_1.text = "Automating Employee Department Updates"
subtitle_1.text = "Bridging the Gap Between HR and Finance using AI\n[Your Name], Accounts Payable Specialist\nAI in Practice - Final Project"

# Slide 2
add_bullet_slide(prs, "Introduction & Background", [
    "My Role: Working in the Finance Department as an Accounts Payable Specialist.",
    "Core Responsibility: Booking personnel costs to the correct departments and cost centers.",
    "The Goal: Ensuring financial reporting is 100% accurate based on employee placement."
])

# Slide 3
add_bullet_slide(prs, "Why I Chose This Topic", [
    "First-hand Frustration: I regularly face issues where employee costs are booked incorrectly.",
    "Investigation: I asked our HR representative about their process.",
    "The Discovery: HR updates department changes manually in BambooHR, but there is no automated system to notify Finance or update our ERP system (NetSuite)."
])

# Slide 4
add_bullet_slide(prs, "The Problem", [
    "The Manual Disconnect: HR successfully updates BambooHR. However, they frequently forget to log into NetSuite to make the exact same change.",
    "The Consequence: Finance relies on NetSuite. Because NetSuite contains outdated information, employee costs get directed to the wrong department.",
    "The Result: Time wasted on manual corrections, inaccurate financial budgets, and data entry fatigue."
])

# Slide 5
add_bullet_slide(prs, "The Proposed Solution", [
    "The Idea: Build an application that automatically updates the employee's department in NetSuite the moment it is changed in BambooHR.",
    "The Method: Using a workflow automation tool to connect the APIs of both systems.",
    "The Trigger: \"When department changes in BambooHR\" -> \"Update record in NetSuite\"."
])

# Slide 6
add_bullet_slide(prs, "The Challenge (Why it's not simple)", [
    "Different Languages: HR and Finance use different terminology.",
    "Example: HR might label a department \"IT Support\", but NetSuite requires a strict financial cost center code, like \"Cost Center 900: IT\".",
    "The Missing Link: A standard integration isn't enough; the data needs to be intelligently translated before it enters NetSuite."
])

# Slide 7
add_bullet_slide(prs, "AI Tools and Techniques", [
    "Integration Platform: Make.com or Zapier to handle the API connections without needing a developer.",
    "The AI Tool: OpenAI (ChatGPT API) integrated directly into the workflow.",
    "The Technique: We use an LLM (Large Language Model) as a \"Translator.\" AI takes the HR department name, compares it against Finance Cost Center rules, and generates the required code."
])

# Slide 8
add_bullet_slide(prs, "How the AI Workflow Looks", [
    "1. Trigger: BambooHR detects a department change.",
    "2. Fetch: The workflow pulls the new HR Department Name.",
    "3. AI Processing: OpenAI receives prompt: \"Map 'Digital Marketing' to correct NetSuite Cost Center.\"",
    "4. AI Output: \"Cost Center 420-Mktg\".",
    "5. Action: The workflow pushes the AI-formatted cost center string into NetSuite."
])

# Slide 9
add_bullet_slide(prs, "Current Stage of the Project", [
    "Status: Planning and Tool Selection Phase.",
    "Completed: Identified the core problem, mapped the manual process, and selected the AI tools (Zapier/Make + OpenAI).",
    "Next Steps: Setting up test accounts (sandboxes) for BambooHR and NetSuite to test the API connection and prompt engineer the OpenAI module."
])

# Slide 10
add_bullet_slide(prs, "Conclusion", [
    "Impact: Eliminates human data-entry error, saves time for both HR and Finance, and guarantees real-time financial accuracy.",
    "The Value of AI: Demonstrates how AI doesn't just write text, but can act as an intelligent \"data translator\" in the background of enterprise systems.",
    "Questions?"
])

# Save the presentation
output_dir = r"c:\Users\kheit\Documents\Xamk\Courses\AI in Practice\Tests"
output_file = os.path.join(output_dir, "Final_Project_Presentation.pptx")
prs.save(output_file)
print(f"Presentation generated successfully: {output_file}")
